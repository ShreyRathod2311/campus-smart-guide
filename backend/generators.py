"""
generators.py — Artifact Generation for CSIS SmartAssist
Generates PPTX slides, Graphviz flowchart PNGs, and AI images on demand.

Image generation priority:
  1. Pollinations.ai  — free, no API key, works immediately
  2. DALL-E 3         — high quality, needs OPENAI_API_KEY (optional)
"""

import os
import uuid
import urllib.parse
from pathlib import Path
from typing import Optional

from dotenv import load_dotenv

load_dotenv()

ARTIFACTS_DIR = Path(os.getenv("ARTIFACTS_DIR", "./artifacts"))
ARTIFACTS_DIR.mkdir(parents=True, exist_ok=True)

# ─────────────────────────────────────────────────────────────────────────────
# BITS Pilani Goa branded color palette
# ─────────────────────────────────────────────────────────────────────────────
BITS_NAVY = (0x1A, 0x2A, 0x5E)       # Navy Blue  — RGB
BITS_GOLD = (0xD4, 0xA0, 0x17)       # Gold       — RGB
BITS_WHITE = (0xFF, 0xFF, 0xFF)
BITS_LIGHT_GREY = (0xF4, 0xF4, 0xF4)


def _rgb(r, g, b):
    """Helper to create pptx RGBColor."""
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    return RGBColor(r, g, b)


# ─────────────────────────────────────────────────────────────────────────────
# PPTX Generation
# ─────────────────────────────────────────────────────────────────────────────

def generate_pptx(title: str, key_points: list[str], filename: Optional[str] = None) -> str:
    """
    Generate a BITS Pilani Goa themed PowerPoint presentation.

    Args:
        title:      Slide deck title (used on cover slide + header)
        key_points: List of bullet point strings (one bullet per point)
        filename:   Output filename (without extension). Auto-generated if None.

    Returns:
        Absolute path to the saved .pptx file.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: Title Slide ──────────────────────────────────────────────────
    blank_layout = prs.slide_layouts[6]  # completely blank
    slide = prs.slides.add_slide(blank_layout)

    # Navy background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*BITS_NAVY)

    # Gold accent bar (left strip)
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0),
        Inches(0.4), Inches(7.5),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(*BITS_GOLD)
    bar.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.5), Inches(11), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.bold = True
    run.font.size = Pt(36)
    run.font.color.rgb = RGBColor(*BITS_WHITE)

    # BITS Pilani Goa subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.7), Inches(4.2), Inches(11), Inches(0.8))
    sf = sub_box.text_frame
    sp = sf.paragraphs[0]
    sp.alignment = PP_ALIGN.LEFT
    sr = sp.add_run()
    sr.text = "CSIS SmartAssist  ·  BITS Pilani Goa"
    sr.font.size = Pt(16)
    sr.font.color.rgb = RGBColor(*BITS_GOLD)

    # ── Slide 2+: Content Slides (6 bullets per slide) ────────────────────────
    BULLETS_PER_SLIDE = 6
    chunks = [key_points[i: i + BULLETS_PER_SLIDE] for i in range(0, len(key_points), BULLETS_PER_SLIDE)]

    for chunk_idx, chunk in enumerate(chunks):
        slide = prs.slides.add_slide(blank_layout)

        # White background
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(*BITS_WHITE)

        # Navy header bar
        header = slide.shapes.add_shape(
            1,
            Inches(0), Inches(0),
            Inches(13.33), Inches(1.1),
        )
        header.fill.solid()
        header.fill.fore_color.rgb = RGBColor(*BITS_NAVY)
        header.line.fill.background()

        # Gold accent left strip
        strip = slide.shapes.add_shape(
            1, Inches(0), Inches(1.1), Inches(0.08), Inches(6.4)
        )
        strip.fill.solid()
        strip.fill.fore_color.rgb = RGBColor(*BITS_GOLD)
        strip.line.fill.background()

        # Slide title in header
        hdr_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(12), Inches(0.8))
        hf = hdr_box.text_frame
        hp = hf.paragraphs[0]
        hr = hp.add_run()
        hr.text = title if chunk_idx == 0 else f"{title} (cont.)"
        hr.font.bold = True
        hr.font.size = Pt(22)
        hr.font.color.rgb = RGBColor(*BITS_WHITE)

        # Content area
        content_box = slide.shapes.add_textbox(Inches(0.4), Inches(1.3), Inches(12.5), Inches(5.8))
        ctf = content_box.text_frame
        ctf.word_wrap = True

        for idx, point in enumerate(chunk):
            para = ctf.paragraphs[idx] if idx == 0 else ctf.add_paragraph()
            para.space_before = Pt(6)
            run = para.add_run()
            run.text = f"• {point}"
            run.font.size = Pt(18)
            run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    # ── Save ──────────────────────────────────────────────────────────────────
    if not filename:
        filename = f"smartassist_{uuid.uuid4().hex[:8]}"
    out_path = ARTIFACTS_DIR / f"{filename}.pptx"
    prs.save(str(out_path))
    print(f"[generators] PPTX saved → {out_path}")
    return str(out_path.resolve())


# ─────────────────────────────────────────────────────────────────────────────
# Graphviz Flowchart Generation
# ─────────────────────────────────────────────────────────────────────────────

def generate_flowchart(steps: list[str], chart_title: str, filename: Optional[str] = None) -> str:
    """
    Generate a directed flowchart PNG using Graphviz.

    Args:
        steps:       Ordered list of process steps.
        chart_title: Title displayed above the flowchart.
        filename:    Output filename (without extension). Auto-generated if None.

    Returns:
        Absolute path to the saved .png file.
    """
    try:
        import graphviz
    except ImportError:
        raise RuntimeError("graphviz package is not installed. Run: pip install graphviz")

    if not filename:
        filename = f"flow_{uuid.uuid4().hex[:8]}"

    # Strip .png if accidentally passed
    filename = filename.replace(".png", "")

    dot = graphviz.Digraph(
        name=chart_title,
        format="png",
        graph_attr={
            "rankdir": "TB",
            "bgcolor": "#FFFFFF",
            "fontname": "Helvetica",
            "label": chart_title,
            "labelloc": "t",
            "fontsize": "20",
            "fontcolor": "#1A2A5E",
            "pad": "0.5",
            "splines": "ortho",
        },
        node_attr={
            "shape": "box",
            "style": "filled,rounded",
            "fillcolor": "#1A2A5E",
            "fontcolor": "#FFFFFF",
            "fontname": "Helvetica",
            "fontsize": "13",
            "margin": "0.3,0.15",
        },
        edge_attr={
            "color": "#D4A017",
            "penwidth": "2",
            "arrowhead": "open",
        },
    )

    # Start / End nodes styled differently
    dot.node("START", "Start", shape="ellipse", fillcolor="#D4A017", fontcolor="#1A2A5E")
    for i, step in enumerate(steps):
        node_id = f"step_{i}"
        dot.node(node_id, step)
    dot.node("END", "End", shape="ellipse", fillcolor="#D4A017", fontcolor="#1A2A5E")

    # Edges
    dot.edge("START", "step_0")
    for i in range(len(steps) - 1):
        dot.edge(f"step_{i}", f"step_{i + 1}")
    dot.edge(f"step_{len(steps) - 1}", "END")

    out_path = ARTIFACTS_DIR / filename
    rendered = dot.render(filename=str(out_path), cleanup=True)
    # graphviz appends .png when format="png"
    final_path = Path(rendered)
    if not final_path.exists():
        # Try with .png appended
        final_path = Path(rendered + ".png") if not rendered.endswith(".png") else final_path

    print(f"[generators] Flowchart saved → {final_path}")
    return str(final_path.resolve())


# ─────────────────────────────────────────────────────────────────────────────
# Intent Detection Helpers
# ─────────────────────────────────────────────────────────────────────────────

def detect_pptx_request(message: str) -> bool:
    """Return True if user is asking for a summary/presentation."""
    keywords = [
        "summary", "presentation", "slides", "pptx", "powerpoint",
        "slide deck", "summarize", "give me a deck", "make a presentation",
    ]
    msg_lower = message.lower()
    return any(k in msg_lower for k in keywords)


def detect_diagram_request(message: str) -> bool:
    """Return True if user is asking for a workflow/flowchart/diagram."""
    keywords = [
        "workflow", "flowchart", "diagram", "flow", "steps for",
        "show me the process", "process for", "visualize", "chart",
    ]
    msg_lower = message.lower()
    return any(k in msg_lower for k in keywords)


def detect_image_request(message: str) -> bool:
    """Return True if user explicitly wants an image/picture generated."""
    keywords = [
        "generate image", "create image", "show image", "picture of",
        "image of", "draw", "generate a picture", "show me a picture",
        "illustration", "generate diagram", "create diagram",
        "show me a diagram", "visual",
    ]
    msg_lower = message.lower()
    return any(k in msg_lower for k in keywords)


def extract_key_points_from_context(context: str, max_points: int = 10) -> list[str]:
    """
    Heuristic extractor: splits context into meaningful bullet points.
    Used when Gemini hasn't produced a JSON key-points list yet.
    """
    lines = [line.strip() for line in context.split("\n") if line.strip()]
    bullets = [
        line.lstrip("•-*").strip()
        for line in lines
        if not line.startswith("**") and not line.startswith("---") and len(line) > 15
    ]
    return bullets[:max_points]


def build_image_prompt(query: str, context: str = "") -> str:
    """
    Build a clean image generation prompt from the user query.
    Strips filler words and adds a professional campus/educational style.
    """
    # Remove command words to get the core subject
    remove_words = [
        "show me", "generate", "create", "draw", "visualize",
        "diagram of", "image of", "picture of", "illustration of",
        "flowchart of", "workflow for", "flow of",
    ]
    subject = query.lower()
    for w in remove_words:
        subject = subject.replace(w, "")
    subject = subject.strip().rstrip("?")

    return (
        f"Professional clean infographic diagram: {subject}. "
        f"BITS Pilani Goa university campus context. "
        f"Navy blue and gold color scheme. White background. "
        f"Modern flat design. No text overlays. High quality."
    )


# ─────────────────────────────────────────────────────────────────────────────
# Free Image Generation — Pollinations.ai (no API key required)
# ─────────────────────────────────────────────────────────────────────────────

# ─────────────────────────────────────────────────────────────────────────────
# Tier 1: Gemini Imagen (uses existing GEMINI_API_KEY, highest quality)
# ─────────────────────────────────────────────────────────────────────────────

async def generate_image_gemini(prompt: str) -> Optional[str]:
    """
    Generate an image using Gemini Imagen 3.
    Requires GEMINI_API_KEY in environment.
    Returns local file path, or None if unavailable.
    """
    gemini_key = os.getenv("GEMINI_API_KEY", "")
    if not gemini_key or gemini_key == "your-gemini-api-key-here":
        return None

    try:
        from google import genai
        from google.genai import types

        client = genai.Client(api_key=gemini_key)
        response = client.models.generate_images(
            model="imagen-3.0-generate-002",
            prompt=prompt,
            config=types.GenerateImagesConfig(
                number_of_images=1,
                aspect_ratio="4:3",
                safety_filter_level="BLOCK_LOW_AND_ABOVE",
            ),
        )
        if not response.generated_images:
            return None

        image_bytes = response.generated_images[0].image.image_bytes
        fname = f"img_{uuid.uuid4().hex[:8]}.png"
        out_path = ARTIFACTS_DIR / fname
        out_path.write_bytes(image_bytes)
        print(f"[generators] Gemini Imagen saved → {out_path} ({len(image_bytes)//1024}KB)")
        return str(out_path.resolve())
    except Exception as e:
        print(f"[generators] Gemini Imagen error: {e}")
        return None


# ─────────────────────────────────────────────────────────────────────────────
# Tier 2: Free Image Generation — Pollinations.ai (no API key)
# ─────────────────────────────────────────────────────────────────────────────

async def generate_image_pollinations(prompt: str) -> Optional[str]:
    """
    Generate an image using Pollinations.ai — completely free, no API key.
    Tries multiple models for resilience.
    Returns local file path, or None on error.
    """
    import httpx

    encoded = urllib.parse.quote(prompt[:500], safe="")
    models = ["flux-schnell", "flux", "turbo"]

    for model in models:
        img_url = (
            f"https://image.pollinations.ai/prompt/{encoded}"
            f"?width=1024&height=768&nologo=true&model={model}"
        )
        try:
            print(f"[generators] Pollinations ({model}) → {img_url[:70]}...")
            async with httpx.AsyncClient(timeout=40.0) as client:
                resp = await client.get(img_url, follow_redirects=True)
                if resp.status_code >= 500:
                    print(f"[generators] Pollinations {model} returned {resp.status_code}, trying next model")
                    continue
                resp.raise_for_status()

                content_type = resp.headers.get("content-type", "")
                if "image" not in content_type:
                    continue

                fname = f"img_{uuid.uuid4().hex[:8]}.png"
                out_path = ARTIFACTS_DIR / fname
                out_path.write_bytes(resp.content)
                print(f"[generators] Pollinations image saved → {out_path} ({len(resp.content)//1024}KB)")
                return str(out_path.resolve())
        except Exception as e:
            print(f"[generators] Pollinations {model} error: {e}")
            continue

    return None


# ─────────────────────────────────────────────────────────────────────────────
# Tier 3: SVG Process Diagram (guaranteed offline fallback)
# ─────────────────────────────────────────────────────────────────────────────

def generate_svg_diagram(steps: list[str], title: str, filename: Optional[str] = None) -> str:
    """
    Generate a clean SVG process-flow diagram (no external tools needed).
    Works 100% offline. Saved as .svg in artifacts dir.
    Returns absolute file path.
    """
    BOX_W, BOX_H = 220, 54
    GAP = 36
    ARROW_H = 30
    PADDING_X = 50
    PADDING_Y = 60
    TITLE_H = 50

    # ─ colours (BITS branding)
    NAVY = "#1A2A5E"
    GOLD = "#D4A017"
    WHITE = "#FFFFFF"
    LIGHT = "#F0F4FF"

    total_h = PADDING_Y + TITLE_H + len(steps) * (BOX_H + ARROW_H) - ARROW_H + PADDING_Y
    total_w = BOX_W + PADDING_X * 2

    lines = []
    lines.append(f'<svg xmlns="http://www.w3.org/2000/svg" width="{total_w}" height="{total_h}" viewBox="0 0 {total_w} {total_h}">')
    lines.append(f'<rect width="{total_w}" height="{total_h}" fill="{LIGHT}"/>')

    # Title
    lines.append(f'<text x="{total_w//2}" y="{PADDING_Y//2+18}" font-family="Arial,sans-serif" font-size="16" font-weight="bold" fill="{NAVY}" text-anchor="middle">{_xml_escape(title)}</text>')
    lines.append(f'<line x1="{PADDING_X}" y1="{PADDING_Y//2+28}" x2="{total_w-PADDING_X}" y2="{PADDING_Y//2+28}" stroke="{GOLD}" stroke-width="2"/>')

    cy = PADDING_Y + TITLE_H
    cx = PADDING_X

    for i, step in enumerate(steps):
        # Box (alternating lightened navy)
        fill = NAVY if i % 2 == 0 else "#2A3D7A"
        lines.append(f'<rect x="{cx}" y="{cy}" width="{BOX_W}" height="{BOX_H}" rx="10" ry="10" fill="{fill}" stroke="{GOLD}" stroke-width="1.5"/>')

        # Step number circle
        lines.append(f'<circle cx="{cx+22}" cy="{cy+BOX_H//2}" r="13" fill="{GOLD}"/>')
        lines.append(f'<text x="{cx+22}" y="{cy+BOX_H//2+5}" font-family="Arial,sans-serif" font-size="12" font-weight="bold" fill="{NAVY}" text-anchor="middle">{i+1}</text>')

        # Step text (truncate if too long)
        label = step[:32] + "…" if len(step) > 32 else step
        lines.append(f'<text x="{cx+44}" y="{cy+BOX_H//2+5}" font-family="Arial,sans-serif" font-size="13" fill="{WHITE}">{_xml_escape(label)}</text>')

        # Arrow (skip after last)
        if i < len(steps) - 1:
            ax = cx + BOX_W // 2
            ay_start = cy + BOX_H
            ay_end = cy + BOX_H + ARROW_H
            lines.append(f'<line x1="{ax}" y1="{ay_start}" x2="{ax}" y2="{ay_end-10}" stroke="{GOLD}" stroke-width="2"/>')
            lines.append(f'<polygon points="{ax},{ay_end} {ax-7},{ay_end-12} {ax+7},{ay_end-12}" fill="{GOLD}"/>')

        cy += BOX_H + ARROW_H

    lines.append('</svg>')
    svg_content = "\n".join(lines)

    if not filename:
        filename = f"diagram_{uuid.uuid4().hex[:8]}"
    out_path = ARTIFACTS_DIR / f"{filename}.svg"
    out_path.write_text(svg_content, encoding="utf-8")
    print(f"[generators] SVG diagram saved → {out_path}")
    return str(out_path.resolve())


def _xml_escape(text: str) -> str:
    """Escape special XML characters."""
    return (
        text.replace("&", "&amp;")
            .replace("<", "&lt;")
            .replace(">", "&gt;")
            .replace('"', "&quot;")
    )


# ─────────────────────────────────────────────────────────────────────────────
# DALL-E Image Generation (optional, needs OPENAI_API_KEY)
# ─────────────────────────────────────────────────────────────────────────────

async def generate_image_dalle(prompt: str) -> Optional[str]:
    """
    Generate an image using DALL-E 3 and save it locally.
    Returns the local file path, or None if OPENAI_API_KEY is not configured.
    """
    import httpx

    openai_key = os.getenv("OPENAI_API_KEY", "")
    if not openai_key:
        return None

    try:
        async with httpx.AsyncClient(timeout=35.0) as client:
            resp = await client.post(
                "https://api.openai.com/v1/images/generations",
                headers={"Authorization": f"Bearer {openai_key}"},
                json={
                    "model": "dall-e-3",
                    "prompt": f"Educational campus diagram, clean, professional: {prompt}",
                    "n": 1,
                    "size": "1024x1024",
                    "quality": "standard",
                },
            )
            resp.raise_for_status()
            img_url = resp.json()["data"][0]["url"]

            img_resp = await client.get(img_url)
            img_resp.raise_for_status()
            fname = f"img_{uuid.uuid4().hex[:8]}.png"
            out_path = ARTIFACTS_DIR / fname
            out_path.write_bytes(img_resp.content)
            print(f"[generators] DALL-E image saved → {out_path}")
            return str(out_path.resolve())
    except Exception as e:
        print(f"[generators] DALL-E error: {e}")
        return None


async def generate_image(prompt: str, steps: Optional[list[str]] = None) -> Optional[str]:
    """
    Smart image generator with 4-tier fallback chain:
      1. Gemini Imagen (if GEMINI_API_KEY is set)
      2. DALL-E 3 (if OPENAI_API_KEY is set)
      3. Pollinations.ai (free, no key)
      4. SVG process diagram (offline, always works)
    """
    # Tier 1: Gemini Imagen
    result = await generate_image_gemini(prompt)
    if result:
        return result

    # Tier 2: DALL-E 3
    result = await generate_image_dalle(prompt)
    if result:
        return result

    # Tier 3: Pollinations.ai
    result = await generate_image_pollinations(prompt)
    if result:
        return result

    # Tier 4: SVG diagram (always succeeds if steps provided)
    print("[generators] All online generators failed, using SVG diagram fallback")
    if steps:
        title = prompt[:60].split(".")[0].strip()
        return generate_svg_diagram(steps, title)
    return None
